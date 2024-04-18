#from data_scraper import reservations
from pptx import Presentation
from data_scraper import get_reservations
from pptx.dml.color import RGBColor  # Import RGBColor explicitly
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from datetime import datetime

#get reservations using the data_scraper
reservations = get_reservations()


template_path = './template.pptx'
prs = Presentation(template_path)


slide = prs.slides[0]

#putting today's date on top of the table
date_shape = slide.shapes[2]
today_date = datetime.now()
today_date_string = today_date.strftime("%A, %Y-%m-%d")
date_shape.text_frame.clear()
p = date_shape.text_frame.add_paragraph()
p.text = today_date_string
p.font.size = Pt(60) 
p.font.name = 'Arial' 
p.alignment = PP_ALIGN.CENTER
p.font.bold = True 

# Loop through shapes in the slide and find the table
for shape in slide.shapes:
    if shape.has_table:
        table = shape.table
        print("found the table!")
        break  # Assuming there's only one table on the slide

# Get the number of rows and columns in the table
num_rows = len(table.rows)
num_cols = len(table.columns)

#dictionary to map the time to a specific cell location
time_to_cell_dict = {"8:30 AM": 1, "9:00 AM": 2, "9:30 AM": 3, "10:00 AM": 4, "10:30 AM": 5, "11:00 AM": 6, "11:30 AM": 7, "12:00 PM": 8,
                    "12:30 PM": 9, "1:00 PM": 10, "1:30 PM": 11, "2:00 PM": 12, "2:30 PM": 13, "3:00 PM": 14, "3:30 PM": 15, "4:00 PM": 16,
                    "4:30 PM": 17, "5:00 PM": 18}

lab_to_cell_dict = {"Lab 1": 1, "Lab 2": 2, "Lab 3": 3}

#loop through the reservations and put them on the table
for reservation in reservations:
    #handle the case where reservations start time or end time is after 5:00 PM - idea use a python representation of hour and say if time is later than 5 pm
    
    if reservation["start_time"] not in time_to_cell_dict:
        top_cell_index = 18
    else:
        top_cell_index = time_to_cell_dict[reservation["start_time"]]

    
    if reservation["end_time"] not in time_to_cell_dict:
        bottom_cell_index = 18
    else:
        bottom_cell_index = time_to_cell_dict[reservation["end_time"]]


    top_cell = table.cell(top_cell_index, lab_to_cell_dict[reservation["lab"]])
    bottom_cell = table.cell(bottom_cell_index, lab_to_cell_dict[reservation["lab"]])

    top_cell.merge(bottom_cell)
    merged_cell = table.cell(top_cell_index, lab_to_cell_dict[reservation["lab"]])

    p = merged_cell.text_frame.add_paragraph()
    p.text = f"{reservation["title"]}\n{reservation["start_time"]} to {reservation["end_time"]}"

    p.font.size = Pt(30) 
    p.font.name = 'Arial' 
    p.font.bold = True  
    p.alignment = PP_ALIGN.CENTER

    #color the cell
    fill = merged_cell.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(244,13,48)

    
prs.save("schedule_of_the_day2.pptx")
